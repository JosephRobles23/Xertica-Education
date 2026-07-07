"""Parser ligero PDF/Office → Markdown (ADR-0006 §5).

MVP sin GPU: usa librerías puras (`pypdf`, `python-docx`, `python-pptx`) con import
perezoso. MinerU queda como swap de adapter en fase 2. `.md`/`.txt` pasan directo.
Instala los extras con `uv sync --extra rag`.
"""
from .base import BaseParserAdapter


class SimpleParserAdapter(BaseParserAdapter):
    async def parse_document(self, file_bytes: bytes, file_name: str) -> str:
        name = file_name.lower()
        if name.endswith((".md", ".txt")):
            return file_bytes.decode("utf-8", errors="replace")
        if name.endswith(".pdf"):
            return self._parse_pdf(file_bytes)
        if name.endswith(".docx"):
            return self._parse_docx(file_bytes)
        if name.endswith(".pptx"):
            return self._parse_pptx(file_bytes)
        raise ValueError(f"Formato no soportado por SimpleParserAdapter: {file_name}")

    @staticmethod
    def _parse_pdf(file_bytes: bytes) -> str:
        import io
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(file_bytes))
        parts = []
        for i, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                parts.append(f"## Página {i}\n\n{text}")
        return "\n\n".join(parts)

    @staticmethod
    def _parse_docx(file_bytes: bytes) -> str:
        import io
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = (para.style.name or "").lower() if para.style else ""
            if style.startswith("heading"):
                lines.append(f"## {text}")
            else:
                lines.append(text)
        return "\n\n".join(lines)

    @staticmethod
    def _parse_pptx(file_bytes: bytes) -> str:
        import io
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if shape.has_text_frame and shape.text.strip()
            ]
            if texts:
                slides.append(f"## Diapositiva {i}\n\n" + "\n\n".join(texts))
        return "\n\n".join(slides)
